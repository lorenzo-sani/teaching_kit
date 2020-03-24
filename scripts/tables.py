from pptx import Presentation

prs = Presentation("files/ppt_demo.pptx")

# text_runs will be populated with a list of strings,
# one for each text run in presentation
#text_runs = []


# slide_num = len(prs.slides)  # to calculate slides number
text_runs = []
text_runs = ['---','\n','layout: presentation','\n','author: '+ prs.core_properties.author,'\n','title: ' + prs.core_properties.title+"_table","\n---"]
n=0
for slide in prs.slides:
    # n=n+1
    for shape in slide.shapes:
        if shape.has_table:
           tab=shape.table
           rows=tab.rows
           n_rows=len(rows)
           col=tab.columns
           n_col=len(col)
           for r in range(n_rows):
               text_runs.append("\n |")
               for c in range(n_col):
                   if r == 1 and c == 0: # this step is needed to add the intermediate row 
                       text_runs.append(":---|---:| ---:| ---:| ---:|\n")
                       text_runs.append(" |"+tab.cell(r,c).text+" |")
                   else:
                       text_runs.append(tab.cell(r,c).text+" |")   
        else:
            continue 

file = open("_presentations/table_"+prs.core_properties.title+".html","w")

for runs in range(len(text_runs)):
    file.write(text_runs[runs]),
file.close()


