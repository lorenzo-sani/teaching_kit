def pptx_extract(path,filename):    
    
    from pptx import Presentation
    import pptx
    from figure_extract import figure_extract
    from table_extract import table_extract
    from text_extract import text_extract

    import os
    from pathlib import Path

    # Read presentation
    prs = Presentation(path+"/"+filename)

    slide_text = []
    n=0
    m=0
    date = '2020-01-01' # Fixed, it's better to keep always the same date for modules in order to avoid problems 

    author = prs.core_properties.author
    title, file_extension = os.path.splitext(filename)
    
    # Assign tags to presentation
    print("- assign tags to the presentation "+title+", press 0 when finished")
    tags=''
    t=''
    while t!="0":
        t = input("type tag: ")
        if t == "0" or t=="":
            continue
        else:
            tags=tags+' '+t
    

    # Initialize text
    slide_text = ['---','\n','categories: module\nexclude: true\nlayout: presentation','\n','author: ', author,'\n','title: ',title,'\n','tags: ',str(tags)]

    Path("_posts/figures/"+title).mkdir(parents=True, exist_ok=True) #check if destination folder for pictures exists and/or creates it

    # Extract and convert
    for slide in prs.slides:
        slide_text.append("\n---\n#")  # new slide, new line, TITLE --- Append is used to add a value at the end of the string
        try: slide_text.append(slide.shapes.title.text)
        except: slide_text.append("")
        
        slide_text.append("\n")
                              
        for shape in slide.shapes:
            if shape.has_text_frame and not shape == slide.shapes.title:
                m=m+1
                text_extract(shape, slide_text, m)
            elif shape == slide.shapes.title:
                 continue
            elif shape.has_table:
                table_extract(shape, slide_text)
            elif shape.has_chart:
                slide_text.append("***MISSING CHART*** insert manually \n")
                    
            elif isinstance(shape, pptx.shapes.placeholder.PlaceholderPicture) or isinstance(shape, pptx.shapes.picture.Picture):
                n=n+1
                figure_extract(shape, slide_text, n, title)
            else:
                slide_text.append("***MISSING OBJECT*** insert manually \n")
                continue

    # Save Slide
    with open("_posts/modules/"+date+"-"+title+".html","w") as presentation_file:
        presentation_file.writelines(slide_text)

