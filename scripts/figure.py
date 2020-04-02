from pptx import Presentation

prs = Presentation("files/ppt_figure.pptx")

prs1= Presentation("files/ppt_figure.pptx")
img = prs1.slides[0].placeholders[1]
graphicframe = prs1.slides[2].placeholders[1]

# text_runs will be populated with a list of strings,
# one for each text run in presentation
#text_runs = []

import pickle
import numpy
import cv2

# pip install opencv-python
#import numpy as np
#import cv2
def save_object(obj, filename):
    with open(filename, 'wb') as output:  # Overwrites any existing file.
        pickle.dump(obj, output, pickle.HIGHEST_PROTOCOL)

# slide_num = len(prs.slides)  # to calculate slides number
slide_text = []

author = prs.core_properties.author
title = prs.core_properties.title
slide_text = ['---','\n','layout: presentation','\n','author: ', author,'\n','title: ',title]
n=0
for slide in prs.slides:
    slide_text.append("\n---\n#")  # new slide, new line, TITLE --- Append is used to add a value at the end of the string
    slide_text.append(slide.shapes.title.text)
    slide_text.append("\n")
    n=n+1
    for shape in slide.shapes:
        if shape.has_text_frame and not shape == slide.shapes.title:
             continue
        elif shape == slide.shapes.title:
            continue
        elif shape.has_table:
           continue
        elif shape.has_chart:
            slide_text.append("***MISSING CHART*** insert manually \n")
            slide_text.append(shape.shape_type)
            #file=open("figures/chart_1.png","w")
            #file.(shape)
            #file.close()
            #fig = shape
            #fig.savefig("figures/chart_1.pdf")
            #fig = shape
            #fig.save("figures/chart.txt")
            #print("***slide number*** ", n)
            #print(fig.name)
            #print(fig.shape_type)
            #print(fig.placeholder_format.element)
            #print(fig.is_placeholder)
            #print(fig.element)
            # PlaceholderGraphicFrame
            #save_object(fig, 'figures/fig_demo.png')
            #plc = shape.placeholder_format
            #print(plc.element)
            #print(plc.idx)
            #print(plc.type)
            

        elif type(shape)==type(img):
            img=shape.image
            name=img.filename
            with open("figures/"+name, 'wb') as presentation_image: 
                 presentation_image.write(img.blob)
         #fig = shape
         #print("***slide number*** ", n)
         #print(fig.name)
         #print(fig.shape_type)
         #print(fig.element)
         #img=fig.image
         #print(img.dpi)
         #print(img.size)
         #print(img.blob) # prints a very long thing
         # image = cv2.imread(img)
         #cv2.imshow('sdfs',img)
         #img.savefig("sdfsd")    


        else:
            slide_text.append("***MISSING OBJECT*** insert manually \n")
                        
            #print(shape.placeholder_format.type.element)
            print(shape.placeholder_format.element)
            print(shape.shape_type)
            continue


#file = open("_presentations/"+prs.core_properties.title+".html","w")

#for runs in range(len(slide_text)):
#    file.write(slide_text[runs]),
#file.close()


#with open("_presentations/"+prs.core_properties.title+"_figures.html","w") as presentation_file:
#    presentation_file.writelines(slide_text)